import os
import logging
import hashlib
from pathlib import Path
from typing import Dict, Union, Set

from pypdf import PdfReader
from docx import Document
from pptx import Presentation

# Assuming VectorStoreManager and generate_embeddings will be imported in the actual execution context
# For this module, we're defining the structure.
# from .vector_store_manager import VectorStoreManager
# from .embedding_generator import generate_embeddings

# Get a logger for this module
logger = logging.getLogger(__name__)

# Supported extensions for document processing
SUPPORTED_EXTENSIONS = [".pdf", ".docx", ".pptx", ".txt", ".md"]


def _calculate_checksum(file_path: Path) -> str:
    """
    Calculates the MD5 checksum of a file.

    Args:
        file_path: Path to the file.

    Returns:
        MD5 checksum as a hex string, or an empty string if file not found or error.
    """
    try:
        hasher = hashlib.md5()
        with open(file_path, 'rb') as f:
            while chunk := f.read(8192): # Read in 8KB chunks
                hasher.update(chunk)
        return hasher.hexdigest()
    except FileNotFoundError:
        logger.error(f"Checksum calculation: File not found at {file_path}")
        return ""
    except Exception as e:
        logger.error(f"Error calculating checksum for {file_path}: {e}")
        return ""


def extract_text_from_md(file_path: Path) -> str:
    """
    Extracts text from a Markdown (MD) file.
    For MD, it's similar to TXT, just reading the content.
    More sophisticated parsing could be added if MD features (like metadata) need specific handling.

    Args:
        file_path: Path to the MD file.

    Returns:
        The extracted text, or an empty string if extraction fails.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        logger.info(f"Successfully extracted text from MD: {file_path}")
        return text
    except Exception as e:
        logger.error(f"Error reading MD file {file_path}: {e}")
        return ""


def extract_text_from_pdf(file_path: Path) -> str:
    """
    Extracts text from a PDF file.

    Args:
        file_path: Path to the PDF file.

    Returns:
        The extracted text, or an empty string if extraction fails.
    """
    text = ""
    try:
        reader = PdfReader(file_path)
        if reader.is_encrypted:
            logger.warning(f"PDF file is encrypted and cannot be read: {file_path}")
            # Attempt to decrypt with an empty password, though this rarely works.
            try:
                reader.decrypt("")
            except Exception as e:
                logger.error(f"Failed to decrypt PDF {file_path}: {e}")
                return ""

        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text
        logger.info(f"Successfully extracted text from PDF: {file_path}")
    except Exception as e:
        logger.error(f"Error extracting text from PDF {file_path}: {e}")
        return ""
    return text

def extract_text_from_docx(file_path: Path) -> str:
    """
    Extracts text from a DOCX file.

    Args:
        file_path: Path to the DOCX file.

    Returns:
        The extracted text, or an empty string if extraction fails.
    """
    text = ""
    try:
        doc = Document(file_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
        logger.info(f"Successfully extracted text from DOCX: {file_path}")
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {file_path}: {e}")
        return ""
    return text.strip()

def extract_text_from_pptx(file_path: Path) -> str:
    """
    Extracts text from a PPTX file.

    Args:
        file_path: Path to the PPTX file.

    Returns:
        The extracted text, or an empty string if extraction fails.
    """
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        logger.info(f"Successfully extracted text from PPTX: {file_path}")
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {file_path}: {e}")
        return ""
    return text.strip()

def extract_text_from_txt(file_path: Path) -> str:
    """
    Extracts text from a TXT file.

    Args:
        file_path: Path to the TXT file.

    Returns:
        The extracted text, or an empty string if extraction fails.
    """
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        logger.info(f"Successfully extracted text from TXT: {file_path}")
        return text
    except Exception as e:
        logger.error(f"Error reading TXT file {file_path}: {e}")
        return ""

def extract_text_from_document(file_path: Path) -> str:
    """
    Determines the file type based on its extension and calls the appropriate
    text extraction function.

    Args:
        file_path: Path to the document.

    Returns:
        The extracted text, or an empty string if the file type is unsupported
        or an error occurs.
    """
    extension = file_path.suffix.lower()
    logger.info(f"Attempting to extract text from: {file_path} (type: {extension})")

    if extension == ".pdf":
        return extract_text_from_pdf(file_path)
    elif extension == ".docx":
        return extract_text_from_docx(file_path)
    elif extension == ".pptx":
        return extract_text_from_pptx(file_path)
    elif extension == ".txt":
        return extract_text_from_txt(file_path)
    elif extension == ".md":
        return extract_text_from_md(file_path)
    else:
        # Option to use unstructured.partition_auto for other types or as a general fallback
        # from unstructured.partition.auto import partition_auto
        # try:
        #     elements = partition_auto(filename=str(file_path))
        #     text = "\n".join([el.text for el in elements])
        #     logger.info(f"Successfully extracted text using unstructured: {file_path}")
        #     return text
        # except Exception as e:
        #     logger.error(f"Error extracting text with unstructured for {file_path}: {e}")
        #     logger.warning(f"Unsupported file type: {extension} for file {file_path}")
        #     return ""
        logger.warning(f"Unsupported file type: {extension} for file {file_path}. Skipping.")
        return ""

def load_documents_from_directory(directory_path: Path) -> Dict[Path, str]:
    """
    Scans a directory (and its subdirectories) for supported document types
    and extracts text from them.

    Args:
        directory_path: Path to the directory to scan.

    Returns:
        A dictionary where keys are file paths (Path objects) and
        values are the extracted text strings.
    """
    extracted_texts: Dict[Path, str] = {}

    if not directory_path.is_dir():
        logger.error(f"Directory not found: {directory_path}")
        return extracted_texts

    logger.info(f"Scanning directory for documents: {directory_path}")
    for file_path in directory_path.rglob("*"): # rglob for recursive globbing
        if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
            logger.info(f"Found supported file: {file_path}")
            text = extract_text_from_document(file_path)
            if text: # Only add if text extraction was successful
                extracted_texts[file_path] = text
        elif file_path.is_file(): # Log only if it's a file but not supported
            logger.debug(f"Skipping unsupported file type: {file_path.suffix} for file {file_path}")

    logger.info(f"Finished scanning directory. Extracted text from {len(extracted_texts)} documents.")
    return extracted_texts


def monitor_knowledge_base(
    kb_path: Path,
    vector_store_manager: 'VectorStoreManager', # Forward reference for typing
    existing_docs_checksums: Dict[str, str]
) -> Dict[str, str]:
    """
    Monitors a knowledge base directory for document changes and updates the vector store.

    Args:
        kb_path: Path to the knowledge base directory.
        vector_store_manager: Instance of VectorStoreManager.
                              Assumes it has add_document(doc_id, content),
                              update_document(doc_id, content), and remove_document(doc_id) methods.
        existing_docs_checksums: A dictionary mapping document path (str) to its checksum.

    Returns:
        The updated existing_docs_checksums dictionary.
    """
    if not kb_path.is_dir():
        logger.error(f"Knowledge base path is not a directory: {kb_path}")
        return existing_docs_checksums

    logger.info(f"Monitoring knowledge base at: {kb_path}")
    current_doc_paths: Set[str] = set()
    updated_checksums = existing_docs_checksums.copy()

    for file_path in kb_path.rglob("*"):
        if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
            doc_id = str(file_path)
            current_doc_paths.add(doc_id)
            new_checksum = _calculate_checksum(file_path)

            if not new_checksum: # Skip if checksum calculation failed
                logger.warning(f"Could not calculate checksum for {file_path}, skipping.")
                continue

            if doc_id not in updated_checksums:
                logger.info(f"New document detected: {file_path}")
                text_content = extract_text_from_document(file_path)
                if text_content:
                    try:
                        # Assuming vector_store_manager.add_document handles chunking and embedding
                        vector_store_manager.add_document(doc_id=doc_id, content=text_content, metadata={"source": doc_id})
                        updated_checksums[doc_id] = new_checksum
                        logger.info(f"Added document to vector store: {doc_id}")
                    except Exception as e:
                        logger.error(f"Failed to add document {doc_id} to vector store: {e}", exc_info=True)
                else:
                    logger.warning(f"No text extracted from new document {file_path}, not adding to vector store.")

            elif updated_checksums[doc_id] != new_checksum:
                logger.info(f"Modified document detected: {file_path}")
                text_content = extract_text_from_document(file_path)
                if text_content:
                    try:
                        # Assuming vector_store_manager.update_document handles re-chunking and re-embedding
                        vector_store_manager.update_document(doc_id=doc_id, content=text_content, metadata={"source": doc_id})
                        updated_checksums[doc_id] = new_checksum
                        logger.info(f"Updated document in vector store: {doc_id}")
                    except Exception as e:
                        logger.error(f"Failed to update document {doc_id} in vector store: {e}", exc_info=True)
                else:
                    logger.warning(f"No text extracted from modified document {file_path}, not updating in vector store.")
            # else: document is unchanged, do nothing

    # Identify and remove deleted documents
    deleted_doc_ids = set(updated_checksums.keys()) - current_doc_paths
    for doc_id_to_remove in deleted_doc_ids:
        logger.info(f"Document removed from knowledge base: {doc_id_to_remove}")
        try:
            vector_store_manager.remove_document(doc_id=doc_id_to_remove)
            del updated_checksums[doc_id_to_remove]
            logger.info(f"Removed document from vector store: {doc_id_to_remove}")
        except Exception as e:
            logger.error(f"Failed to remove document {doc_id_to_remove} from vector store: {e}", exc_info=True)

    logger.info("Knowledge base monitoring cycle complete.")
    return updated_checksums


if __name__ == '__main__':
    # This section is for basic testing.
    # Proper testing requires a VectorStoreManager mock or instance and actual files.
    logging.basicConfig(
        level=logging.DEBUG,
        format='%(asctime)s - %(name)s - %(levelname)s - [%(filename)s:%(lineno)d] - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )
    logger.info("Document loader module - testing monitor_knowledge_base (conceptual)")

    # --- Setup a mock VectorStoreManager for testing ---
    class MockVectorStoreManager:
        def __init__(self):
            self.documents = {}
            logger.debug("MockVectorStoreManager initialized.")

        def add_document(self, doc_id: str, content: str, metadata: Dict = None):
            logger.info(f"[MockVSM] ADDING: {doc_id}, content length: {len(content)}, metadata: {metadata}")
            if doc_id in self.documents:
                logger.warning(f"[MockVSM] Document {doc_id} already exists. Overwriting for add_document.")
            self.documents[doc_id] = {"content": content, "metadata": metadata or {}}

        def update_document(self, doc_id: str, content: str, metadata: Dict = None):
            logger.info(f"[MockVSM] UPDATING: {doc_id}, content length: {len(content)}, metadata: {metadata}")
            if doc_id not in self.documents:
                logger.warning(f"[MockVSM] Document {doc_id} not found for update. Adding it instead.")
            self.documents[doc_id] = {"content": content, "metadata": metadata or {}}

        def remove_document(self, doc_id: str):
            logger.info(f"[MockVSM] REMOVING: {doc_id}")
            if doc_id in self.documents:
                del self.documents[doc_id]
            else:
                logger.warning(f"[MockVSM] Document {doc_id} not found for removal.")

        def get_document_count(self):
            return len(self.documents)

    # --- Setup a temporary test knowledge base directory ---
    test_kb_dir = Path("./temp_kb_for_testing")
    test_kb_dir.mkdir(parents=True, exist_ok=True)

    mock_vsm = MockVectorStoreManager()
    current_checksums = {}

    # --- Test Cycle 1: Add new files ---
    logger.info("\n--- Test Cycle 1: Adding new files ---")
    with open(test_kb_dir / "doc1.txt", "w") as f: f.write("Hello world from doc1")
    with open(test_kb_dir / "doc2.md", "w") as f: f.write("# Markdown Content\nTest for doc2.")
    # Unsupported file
    with open(test_kb_dir / "doc.unsupported", "w") as f: f.write("Unsupported content")


    current_checksums = monitor_knowledge_base(test_kb_dir, mock_vsm, current_checksums)
    assert str(test_kb_dir / "doc1.txt") in current_checksums
    assert str(test_kb_dir / "doc2.md") in current_checksums
    assert mock_vsm.get_document_count() == 2
    logger.info(f"Checksums after cycle 1: {current_checksums}")
    logger.info(f"MockVSM docs after cycle 1: {mock_vsm.documents.keys()}")


    # --- Test Cycle 2: Modify one file, add another, no deletions ---
    logger.info("\n--- Test Cycle 2: Modifying and adding files ---")
    with open(test_kb_dir / "doc1.txt", "w") as f: f.write("Hello world from doc1 - updated") # Modified
    with open(test_kb_dir / "doc3.pdf", "w") as f: f.write("PDF content (simulated for txt)") # New PDF (will use txt extractor)

    # Simulate a delay or ensure modification times are different if checksums were based on mtime
    # For content-based checksums, this is not strictly necessary but good practice.
    import time
    time.sleep(0.1)

    checksum_before_update = current_checksums[str(test_kb_dir / "doc1.txt")]
    current_checksums = monitor_knowledge_base(test_kb_dir, mock_vsm, current_checksums)
    assert current_checksums[str(test_kb_dir / "doc1.txt")] != checksum_before_update
    assert str(test_kb_dir / "doc3.pdf") in current_checksums # .pdf is supported, uses txt extractor for this basic test file
    assert mock_vsm.get_document_count() == 3 # doc1 (updated), doc2 (same), doc3 (new)
    logger.info(f"Checksums after cycle 2: {current_checksums}")
    logger.info(f"MockVSM docs after cycle 2: {mock_vsm.documents.keys()}")

    # --- Test Cycle 3: Delete one file, one remains unchanged ---
    logger.info("\n--- Test Cycle 3: Deleting a file ---")
    (test_kb_dir / "doc2.md").unlink() # Delete doc2.md

    current_checksums = monitor_knowledge_base(test_kb_dir, mock_vsm, current_checksums)
    assert str(test_kb_dir / "doc2.md") not in current_checksums
    assert mock_vsm.get_document_count() == 2 # doc1, doc3 remain
    logger.info(f"Checksums after cycle 3: {current_checksums}")
    logger.info(f"MockVSM docs after cycle 3: {mock_vsm.documents.keys()}")

    # --- Test Cycle 4: Emptying directory ---
    logger.info("\n--- Test Cycle 4: Emptying directory ---")
    for f_path in test_kb_dir.glob("*"):
        f_path.unlink()

    current_checksums = monitor_knowledge_base(test_kb_dir, mock_vsm, current_checksums)
    assert not current_checksums # Checksums should be empty
    assert mock_vsm.get_document_count() == 0 # All documents removed
    logger.info(f"Checksums after cycle 4: {current_checksums}")
    logger.info(f"MockVSM docs after cycle 4: {mock_vsm.documents.keys()}")


    # --- Clean up ---
    logger.info("\n--- Cleaning up test directory ---")
    # (test_kb_dir / "doc1.txt").unlink(missing_ok=True) # Already deleted in cycle 4
    # (test_kb_dir / "doc3.pdf").unlink(missing_ok=True) # Already deleted in cycle 4
    test_kb_dir.rmdir()
    logger.info(f"Test directory {test_kb_dir} removed.")

    logger.info("Document loader module - monitor_knowledge_base test run complete.")
    # This section can be used for basic testing of the module later
    # Note: For this __main__ block to work standalone with proper logging,
    # you'd need to add a basicConfig here, but it's generally not done for library modules.
    # The application using this module (like ingest_documents.py) should set up logging.
    logger.info("Document loader module can be tested here if logging is configured by the caller.")
    # Example:
    # logging.basicConfig(level=logging.DEBUG, format='%(asctime)s - %(name)s - %(levelname)s - [%(filename)s:%(lineno)d] - %(message)s')
    # test_doc_path = Path("path_to_your_test_document.pdf")
    # if test_doc_path.exists():
    #     text = extract_text_from_document(test_doc_path)
    #     logger.info(f"Extracted text (first 100 chars): {text[:100]}")
    # else:
    #     logger.warning(f"Test document not found: {test_doc_path}")
    pass
